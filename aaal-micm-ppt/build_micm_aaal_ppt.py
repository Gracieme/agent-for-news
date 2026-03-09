#!/usr/bin/env python3
"""
MICM · AAAL 2026 — Designer-grade deck.
Editorial split cover, full-bleed section dividers, key-number slide, minimal tables.
+ Slide transitions (fade) and timing via OOXML post-process.
AAAL-compliant: Arial, 24pt+ body, high contrast.
"""
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Refined palette
BG = RGBColor(0xFB, 0xF9, 0xF6)          # warm paper
TEXT = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_LIGHT = RGBColor(0x5C, 0x5C, 0x5C)
TEAL = RGBColor(0x0D, 0x5C, 0x5C)        # rich teal
TEAL_SOFT = RGBColor(0xE6, 0xF0, 0xF0)   # soft fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"
W = 10.0
H = 5.625


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _p(para, size_pt, bold=False, color=None, align=None, space_after=8):
    para.font.name = FONT
    para.font.size = Pt(size_pt)
    para.font.bold = bold
    if color:
        para.font.color.rgb = color
    if align is not None:
        para.alignment = align
    para.space_after = Pt(space_after)


# ---- COVER: editorial split (teal left 36%, content right) ----
def add_cover(prs, main_title, sub_title, byline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    # Left block: full-height teal (36% of 10" = 3.6")
    block_w = Inches(W * 0.36)
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, block_w, prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = TEAL
    left.line.fill.background()
    # "AAAL 2026" on teal
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.7), Inches(1.2), Inches(0.4))
    tb.text_frame.paragraphs[0].text = "AAAL 2026"
    _p(tb.text_frame.paragraphs[0], 18, color=WHITE)
    tb.text_frame.paragraphs[0].font.underline = True
    # Content on cream
    cx = Inches(0.5 + W * 0.36)
    cw = Inches(W - 0.6 - W * 0.36)
    t2 = slide.shapes.add_textbox(cx, Inches(1.15), cw, Inches(2.2))
    tf = t2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_title
    _p(p, 34, bold=True, color=TEXT, space_after=16)
    p2 = tf.add_paragraph()
    p2.text = sub_title
    _p(p2, 22, color=TEXT_LIGHT, space_after=28)
    p3 = tf.add_paragraph()
    p3.text = byline
    _p(p3, 20, color=TEXT_LIGHT)
    return slide


# ---- SECTION DIVIDER: full teal, big number + title ----
def add_part(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, TEAL)
    # Big number
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(3), Inches(1.2))
    box.text_frame.paragraphs[0].text = number
    _p(box.text_frame.paragraphs[0], 96, bold=True, color=WHITE)
    box.text_frame.paragraphs[0].font.name = FONT
    # Title
    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(6), Inches(0.8))
    t2.text_frame.paragraphs[0].text = title
    _p(t2.text_frame.paragraphs[0], 32, bold=True, color=WHITE)
    return slide


# ---- KEY NUMBER: one big stat ----
def add_key_number(prs, big_text, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    # Centered big number/quote
    box = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(8), Inches(1.4))
    tf = box.text_frame
    tf.paragraphs[0].text = big_text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _p(tf.paragraphs[0], 72, bold=True, color=TEAL)
    p2 = tf.add_paragraph()
    p2.text = caption
    p2.alignment = PP_ALIGN.CENTER
    _p(p2, 24, color=TEXT_LIGHT, space_after=0)
    p2.space_before = Pt(12)
    return slide


# ---- CONTENT: left block + title + body ----
BLOCK_W = Inches(0.38)
CONTENT_X = Inches(0.58)
CONTENT_W = Inches(8.5)
TOP_Y = Inches(0.45)
BODY_Y = Inches(1.35)
BODY_H = Inches(3.85)


def add_content(prs, title_text, bullets, pres):
    slide = prs.slides.add_slide(pres.slide_layouts[6])
    _bg(slide)
    # Left teal block
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BLOCK_W, pres.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    # Title
    th = slide.shapes.add_textbox(CONTENT_X, TOP_Y, CONTENT_W, Inches(0.85))
    th.text_frame.paragraphs[0].text = title_text
    _p(th.text_frame.paragraphs[0], 28, bold=True, color=TEAL)
    # Body
    body = slide.shapes.add_textbox(CONTENT_X, BODY_Y, CONTENT_W, BODY_H)
    btf = body.text_frame
    btf.word_wrap = True
    for i, line in enumerate(bullets[:6]):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        _p(p, 24, color=TEXT, space_after=14)
    return slide


# ---- TABLE: minimal header, clean rows ----
def add_table(prs, title_text, headers, rows, col_widths, pres):
    slide = prs.slides.add_slide(pres.slide_layouts[6])
    _bg(slide)
    # Left block
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BLOCK_W, pres.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    # Title
    th = slide.shapes.add_textbox(CONTENT_X, TOP_Y, CONTENT_W, Inches(0.6))
    th.text_frame.paragraphs[0].text = title_text
    _p(th.text_frame.paragraphs[0], 26, bold=True, color=TEAL)
    # Table
    nrows, ncols = 1 + len(rows), len(headers)
    row_h = Inches(0.42)
    tw = Inches(8.2)
    table = slide.shapes.add_table(nrows, ncols, Inches(0.55), Inches(1.2), tw, row_h * nrows).table
    for c, w in enumerate(col_widths):
        if c < ncols:
            table.columns[c].width = Inches(w)
    for c in range(ncols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        cell.text = headers[c]
        for p in cell.text_frame.paragraphs:
            p.font.name, p.font.size = FONT, Pt(22)
            p.font.bold, p.font.color.rgb = True, WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c >= ncols:
                break
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.name, p.font.size = FONT, Pt(22)
                p.font.color.rgb = TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


# ---- ANIMATION / TRANSITION: OOXML post-process ----
# python-pptx has no animation API; inject slide transition (fade) into saved file
TRANSITION_FRAGMENT = '  <p:transition spd="med"><p:fade thruBlk="0"/></p:transition>\n'


def _inject_transition_into_slide_xml(xml_bytes):
    """Insert p:transition (fade) before closing </p:sld>. Returns bytes."""
    text = xml_bytes.decode("utf-8")
    if "<p:transition" in text:
        return xml_bytes
    # Insert once before the slide root closing tag
    new_text = text.replace("</p:sld>", TRANSITION_FRAGMENT + "</p:sld>", 1)
    return new_text.encode("utf-8")


def add_slide_transitions(pptx_path):
    """Add fade transition (medium speed) to every slide in the saved .pptx."""
    path = Path(pptx_path)
    if not path.exists():
        return
    buf = BytesIO()
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    data = _inject_transition_into_slide_xml(data)
                zout.writestr(name, data)
    buf.seek(0)
    with open(path, "wb") as f:
        f.write(buf.read())
    print("  + Slide transitions (fade, medium) applied to all slides.")


# ---- CLOSING: full teal ----
def add_closing(prs, title_line, lines_below):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, TEAL)
    box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(2))
    tf = box.text_frame
    tf.paragraphs[0].text = title_line
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _p(tf.paragraphs[0], 48, bold=True, color=WHITE)
    for line in lines_below:
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        _p(p, 24, color=WHITE)
        p.space_before = Pt(10)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # Cover
    add_cover(prs,
        "Translanguaging as Mediational Architecture",
        "Reconceptualizing negotiation, reflection, and interactional competence in CSL peer interaction",
        "Gracie Jiaxin Shen · University of Massachusetts Amherst"
    )

    add_part(prs, "01", "The puzzle")
    add_content(prs, "Research focus", [
        "NoM→LRE rate drops at novice-high (20.2%) then rises at intermediate (79.2%).",
        "MICM: translanguaging as primary mediational architecture—revision of IH automaticity assumption. Evidence: 11 learners, 21 episodes.",
    ], prs)
    add_content(prs, "The assumption we revise", [
        "Long (1996): negotiation → optimal SLA. Assumption: negotiation and metalinguistic reflection linked by an automatic mechanism in the target language.",
        "Data: Lexical Explanation → LRE 100% (C126) vs 14.3% (C246). Mechanism depends on translanguaging.",
    ], prs)

    add_key_number(prs, "100%  →  14.3%", "Lexical Explanation → LRE · C126 vs C246")
    add_content(prs, "MICM central claim", [
        "Translanguaging = primary mediational architecture: converts communicative uncertainty into metalinguistic reflection.",
        "In CSL, metalinguistic Chinese is insufficient for that precision; cross-linguistic resources provide the medium. (Hypothesis for exploratory scale.)",
    ], prs)

    add_part(prs, "02", "Theory & literature")
    add_content(prs, "Interaction Hypothesis under revision", [
        "Long (1996): negotiation → modified input, pushed output, noticing. Assumption: optimal in target-language-only (challenged but not reformulated).",
        "MICM: translanguaging is the mechanism that converts—or does not convert—negotiation into metalinguistic reflection. Revision, not extension.",
    ], prs)
    add_content(prs, "LREs & translanguaging", [
        "LREs: CSL expansion → tonal (pitch contour → meaning), orthographic (radical semantics).",
        "Translanguaging (García & Wei 2014; Li Wei 2018): MICM repositions as independent mediational variable.",
    ], prs)
    add_content(prs, "MICM Claim 1: Mediational architecture", [
        "Scaffolding = temporary support; mediational architecture (Vygotsky) = function organized through the resource. Example: 又 sounds better — precision unavailable in 对/不对.",
        "Skehan (1998) trade-off: at novice-high, TL predominantly functional → bottleneck.",
    ], prs)
    add_table(prs, "MICM Claims 2 & 3",
        ["Claim", "Summary"],
        [
            ["Claim 2", "Capacity to deploy translanguaging metalinguistically may be integral to IC in typologically distant L2 contexts—constitutively multilingual (hypothesis)"],
            ["Claim 3", "Non-linear pattern: novice-high L2 expansion → TL predominantly functional → bottleneck. Resolves at intermediate (79.2%)"],
        ], [1.4, 6.8], prs)

    add_part(prs, "03", "Method & results")
    add_table(prs, "Context & participants",
        ["Course", "Level", "Episodes", "N"],
        [
            ["CHN 126", "Novice-low", "9", "—"],
            ["CHN 246", "Novice-high", "6", "—"],
            ["CHN 326", "Intermediate-low", "6", "—"],
            ["Total", "—", "21", "11 undergrads"],
        ], [2.2, 2.0, 1.5, 2.5], prs)
    add_table(prs, "Tasks & coding",
        ["Component", "Operationalization"],
        [
            ["Tasks", "Information-gap, decision-making, cultural discussion; ~10 h peer interaction, Jefferson transcription"],
            ["NoM", "Clarification requests, confirmation checks, help requests, self-correction, reformulations"],
            ["LREs", "Lexical, grammatical, phonological/tonal, orthographic (CSL-expanded)"],
            ["Translanguaging", "vs code-switching by functional + sequential criteria (κ = .85)"],
        ], [2.0, 6.2], prs)
    add_table(prs, "CSL-specific LRE type distributions",
        ["LRE type", "C126 NL", "C246 NH", "C326 IL", "Total"],
        [
            ["Grammatical", "2 (22.2%)", "1 (5.6%)", "25 (64.1%)", "28"],
            ["Lexical", "7 (77.8%)", "14 (77.8%)", "12 (30.8%)", "33"],
            ["Phonological/Tonal", "0", "3 (16.7%)", "2 (5.1%)", "5"],
            ["Total", "9", "18", "39", "66"],
        ], [1.6, 1.5, 1.5, 1.5, 0.9], prs)
    add_table(prs, "NoM strategy distributions (χ²=26.67, p=.009)",
        ["Strategy", "C126 (n=86)", "C246 (n=89)", "C326 (n=53)"],
        [
            ["Question generation", "17 (19.8%)", "24 (27.0%)", "22 (41.5%)"],
            ["Confirmation check", "27 (31.4%)", "15 (16.9%)", "16 (30.2%)"],
            ["Self-correction", "17 (19.8%)", "18 (20.2%)", "8 (15.1%)"],
            ["Lexical explanation", "7 (8.1%)", "14 (15.7%)", "1 (1.9%)"],
        ], [2.4, 2.0, 2.0, 2.0], prs)
    add_table(prs, "NoM→LRE transition rates (bottleneck)",
        ["Group", "Total NoM", "Total LREs", "NoM→LRE", "Rate"],
        [
            ["C126 Novice-Low (9 ep.)", "88", "9", "26", "29.5%"],
            ["C246 Novice-High (6 ep.)", "89", "18", "18", "20.2% ↓"],
            ["C326 Interm.-Low (6 ep.)", "53", "39", "42", "79.2% ↑"],
        ], [2.3, 1.4, 1.4, 1.6, 1.2], prs)
    add_table(prs, "Bottleneck fingerprint: Lexical Explanation → LRE",
        ["Level", "Rate", "Interpretation"],
        [
            ["C126 Novice-low", "100%", "Every vocabulary provision triggered metalinguistic reflection"],
            ["C246 Novice-high", "14.3%", "Vocabulary resolves gap; metalinguistic opportunity not taken"],
        ], [2.0, 1.6, 4.6], prs)
    add_table(prs, "Qualitative: four cases",
        ["Case", "Level", "Summary"],
        [
            ["1", "C246", "A: structural uncertainty (比…更好?); B topic-shifts. TL affective, not metalinguistic"],
            ["2", "C246", "A: 很方便 (→ 更方便); B unremarked. Fossilization concern"],
            ["3", "C246", "A: ‘or no the sun’ → 晴天; B provides 晴天. Functional TL; comparative error unremarked"],
            ["4", "C326", "B: 又 sounds better. LRE with TL as metalinguistic medium; bottleneck resolution"],
        ], [0.6, 1.0, 6.6], prs)

    add_part(prs, "04", "Discussion & implications")
    add_content(prs, "Revising the Interaction Hypothesis", [
        "More negotiation ≠ more metalinguistic development at C246. MICM: mechanism is translanguaging-mediated.",
        "Multilingual Interaction Hypothesis: conditions under which Long’s account is most fully operative.",
    ], prs)
    add_content(prs, "Fossilization risk", [
        "Comparative errors (很方便→更方便; 比冬天很暖和→更暖和) missed at novice-high → fossilization candidates (Selinker 1972). Testable: longitudinal C246→C326.",
    ], prs)
    add_table(prs, "Implications for pedagogy",
        ["Focus", "Implication"],
        [
            ["Novice-high", "Scaffold metalinguistic reflection (e.g. structured post-negotiation reflection phases)"],
            ["Translanguaging", "Legitimize functional TL at novice; cultivate metalinguistic TL at intermediate"],
            ["Task design", "Sequence NoM-eliciting, LRE-eliciting, and translanguaging-eliciting tasks within units"],
        ], [1.8, 6.4], prs)
    add_table(prs, "Conclusion: three claims",
        ["Claim", "Summary"],
        [
            ["1", "Translanguaging as primary mediational architecture in these sequences"],
            ["2", "IC in typologically distant L2 may be constitutively multilingual"],
            ["3", "Development non-linear in a theoretically predicted way (bottleneck then resolution)"],
        ], [1.0, 7.2], prs)

    add_closing(prs, "Thank you", [
        "Gracie Jiaxin Shen",
        "jiaxinshen1208@gmail.com",
        "University of Massachusetts Amherst",
    ])

    out = "/Users/geleixizhibao/Downloads/tarot-app/aaal-micm-ppt/MICM_AAAL2026.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    add_slide_transitions(out)


if __name__ == "__main__":
    main()
